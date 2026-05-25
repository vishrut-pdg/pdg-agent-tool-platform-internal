import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from onyx.file_processing.extract_file_text import extract_pptx_images
from onyx.file_processing.extract_file_text import pptx_to_text
from onyx.file_processing.extract_file_text import read_pptx_file


def _make_1x1_png() -> bytes:
    """Minimal valid 1x1 white PNG (67 bytes)."""
    import struct
    import zlib

    signature = b"\x89PNG\r\n\x1a\n"

    def _chunk(ctype: bytes, data: bytes) -> bytes:
        crc = zlib.crc32(ctype + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + ctype + data + struct.pack(">I", crc)

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    raw_row = b"\x00\xff\xff\xff"
    idat = zlib.compress(raw_row)
    return (
        signature + _chunk(b"IHDR", ihdr) + _chunk(b"IDAT", idat) + _chunk(b"IEND", b"")
    )


def _make_pptx_with_image() -> io.BytesIO:
    """Create an in-memory pptx with a text slide and an embedded PNG image."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add a title via a textbox
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(6), Inches(1))
    txBox.text_frame.text = "Slide with image"

    # Embed a real PNG image
    png_bytes = _make_1x1_png()
    slide.shapes.add_picture(
        io.BytesIO(png_bytes), Inches(1), Inches(2), Inches(2), Inches(2)
    )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _make_pptx_with_chart() -> io.BytesIO:
    """Create an in-memory pptx with one text slide and one chart slide."""
    prs = Presentation()

    # Slide 1: text only
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Introduction"
    slide1.placeholders[1].text = "This is the first slide."

    # Slide 2: chart
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3"]
    chart_data.add_series("Revenue", (100, 200, 300))
    slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(4),
        chart_data,
    )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _make_pptx_without_chart() -> io.BytesIO:
    """Create an in-memory pptx with a single text-only slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Hello World"
    slide.placeholders[1].text = "Some content here."

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


class TestPptxToText:
    def test_chart_is_omitted(self) -> None:
        # Precondition
        pptx_file = _make_pptx_with_chart()

        # Under test
        result = pptx_to_text(pptx_file)

        # Postcondition
        assert "Introduction" in result
        assert "first slide" in result
        assert "[chart omitted]" in result
        # The actual chart data should NOT appear in the output.
        assert "Revenue" not in result
        assert "Q1" not in result

    def test_text_only_pptx(self) -> None:
        # Precondition
        pptx_file = _make_pptx_without_chart()

        # Under test
        result = pptx_to_text(pptx_file)

        # Postcondition
        assert "Hello World" in result
        assert "Some content" in result
        assert "[chart omitted]" not in result


class TestExtractPptxImages:
    def test_extracts_embedded_image(self) -> None:
        pptx_file = _make_pptx_with_image()

        images = list(extract_pptx_images(pptx_file))

        assert len(images) == 1
        img_bytes, img_name = images[0]
        assert len(img_bytes) > 0
        assert img_bytes[:4] == b"\x89PNG"
        assert img_name  # non-empty name

    def test_no_images_in_text_only_pptx(self) -> None:
        pptx_file = _make_pptx_without_chart()

        images = list(extract_pptx_images(pptx_file))

        assert images == []

    def test_invalid_file_yields_nothing(self) -> None:
        bad_file = io.BytesIO(b"not a zip file")

        images = list(extract_pptx_images(bad_file))

        assert images == []


class TestReadPptxFile:
    def test_returns_text_without_images_by_default(self) -> None:
        pptx_file = _make_pptx_with_image()

        text, images = read_pptx_file(pptx_file, file_name="test.pptx")

        assert "Slide with image" in text
        assert images == []

    def test_extract_images_returns_list(self) -> None:
        pptx_file = _make_pptx_with_image()

        text, images = read_pptx_file(
            pptx_file, file_name="test.pptx", extract_images=True
        )

        assert "Slide with image" in text
        assert len(images) == 1
        img_bytes, img_name = images[0]
        assert img_bytes[:4] == b"\x89PNG"
        assert img_name

    def test_callback_streams_instead_of_collecting(self) -> None:
        pptx_file = _make_pptx_with_image()
        collected: list[tuple[bytes, str]] = []

        def callback(data: bytes, name: str) -> None:
            collected.append((data, name))

        text, images = read_pptx_file(
            pptx_file,
            file_name="test.pptx",
            extract_images=True,
            image_callback=callback,
        )

        assert "Slide with image" in text
        assert len(collected) == 1
        assert collected[0][0][:4] == b"\x89PNG"
        # Returned list is empty when callback is used
        assert images == []
